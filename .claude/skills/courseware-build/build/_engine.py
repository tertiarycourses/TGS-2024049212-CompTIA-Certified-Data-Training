#!/usr/bin/env python3
"""Generate the Gemini Agent ADK course slide deck (all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py so the deck stays 100% aligned with the LP,
LG and labs.
"""
import os, sys, copy, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4 + DOMAIN5

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4,
        font="Arial",wrap=True):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name=font
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

# ---------------- motion: restrained transitions + build animation ----------------
# House rule: ONE transition family for the whole deck (morph is unavailable in the
# OOXML PowerPoint 2010 transition set that python-pptx can emit, so we use a short
# push/fade pair). Content slides fade; section dividers push. Nothing else moves.
P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"

def _transition(s, kind="fade", speed="med"):
    """Attach a restrained slide transition. kind: fade | push | wipe."""
    sld = s._element
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", speed)
    tr.set("advClick", "1")
    if kind == "fade":
        etree.SubElement(tr, qn("p:fade"))
    elif kind == "push":
        el = etree.SubElement(tr, qn("p:push")); el.set("dir", "l")
    elif kind == "wipe":
        el = etree.SubElement(tr, qn("p:wipe")); el.set("dir", "r")
    # keep the transition last in the slide element (schema order)
    sld.append(tr)
    return tr

# Appear-on-click build for a list of shape ids — used ONLY on process maps so the
# trainer can reveal one stage at a time. No spins, no flying, no sound.
_TIMING = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
<p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">
<p:childTnLst>{pars}</p:childTnLst></p:cTn>
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""

_PAR = """<p:par><p:cTn id="{i0}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="{i1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="{i2}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="{nt}">
<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
<p:set><p:cBhvr><p:cTn id="{i3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{i4}" dur="400"/>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>
</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>"""

def _build_on_click(s, spids):
    """Fade each shape id in on click, in order. Used sparingly (process maps)."""
    if not spids:
        return
    pars, nid = [], 10
    for k, spid in enumerate(spids):
        pars.append(_PAR.format(i0=nid, i1=nid+1, i2=nid+2, i3=nid+3, i4=nid+4,
                                spid=spid, nt="clickEffect" if k == 0 else "afterEffect"))
        nid += 10
    sld = s._element
    for old in sld.findall(qn("p:timing")):
        sld.remove(old)
    sld.append(etree.fromstring(_TIMING.format(pars="".join(pars))))

def connector(s, x1, y1, x2, y2, color, width=Pt(2.0), arrow=True):
    """A REAL PowerPoint connector line (not a typed arrow glyph)."""
    cx, cy = min(x1, x2), min(y1, y2)
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = width
    ln = cn.line._get_or_add_ln()
    if arrow:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return cn

def chevron(s, x, y, w, h, color):
    """A staged chevron shape — the real CHEVRON autoshape, not a text glyph."""
    sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def roundrect(s, x, y, w, h, color, line=None, adj=0.10):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = adj
    except Exception: pass
    return sp

def diamond(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp

def label_in(sp, text, size, color, bold=True):
    """Put centred text inside an autoshape."""
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Arial"
    return sp

PAGE={"n":1}   # the cover is slide 1 and carries no number, so numbering starts at 2
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def _ellipsis(text,limit):
    """Truncate on a WORD boundary with an ellipsis — never slice mid-word."""
    t=" ".join(str(text).split())
    if len(t)<=limit: return t
    cut=t[:limit]
    sp=cut.rfind(" ")
    if sp>limit*0.55: cut=cut[:sp]
    return cut.rstrip(" ,.;:-—(") + "…"

def _short_cmd(cmd,limit=30):
    """Shorten a shell command for a caption. Paths/URLs have no spaces, so a
    word-boundary ellipsis cuts them mid-token — drop the runner prefix and keep the
    meaningful tail (the script/target) instead."""
    c=" ".join(str(cmd).split())
    for pre in ("uv run python ","uv run ","python3 ","python ","bash ","sh "):
        if c.startswith(pre): c=c[len(pre):]; break
    if len(c)<=limit: return c
    parts=c.split(" ")
    head=parts[0]
    rest=" ".join(parts[1:]).strip()
    if rest:
        # a URL argument → keep the verb and the host, drop the path
        m=re.match(r'https?://([^/\s]+)',rest)
        if m:
            for cand in (f"{head} {m.group(1)}/…", f"{head} {m.group(1)}"):
                if len(cand)<=limit: return cand
        # "git clone …/repo" — keep the verb plus the final path segment
        tail=rest.rstrip("/").split("/")[-1]
        cand=f"{head} …/{tail}" if "/" in rest else f"{head} {tail}"
        if len(cand)<=limit: return cand
        if len(head)+2<=limit: return _ellipsis(head,limit-2)+" …"
    if "/" in c:                   # a bare path → keep the last segment
        tail=c.rstrip("/").split("/")[-1]
        if len(tail)+2<=limit: return "…/"+tail
    return _ellipsis(c,limit)

def _wrap_cmd(cmd, per_line=104, max_lines=3):
    """Break a long command across lines on WHITESPACE so every token stays whole.
    A token longer than per_line (a URL or a long SQL literal) is hard-split only as a
    last resort. Returns a list of lines; the final line is ellipsised if it still
    overflows max_lines."""
    words=cmd.split(" ")
    lines=[]; cur=""
    for w in words:
        cand=w if not cur else cur+" "+w
        if len(cand)<=per_line:
            cur=cand; continue
        if cur: lines.append(cur); cur=""
        while len(w)>per_line:
            # Prefer a SYNTACTIC break (; or , or a closing paren) over a fixed column,
            # so a python3 -c "..." one-liner never splits in the middle of a token.
            window=w[:per_line]
            # Only break at a separator that sits OUTSIDE a quoted literal — a break
            # inside a regex or string renders as a space and changes its meaning.
            cut=-1
            depth_s=depth_d=0
            for k,ch in enumerate(window):
                if ch=="'" and depth_d==0: depth_s^=1
                elif ch=='"' and depth_s==0: depth_d^=1
                elif ch in ";,)" and depth_s==0 and depth_d==0: cut=k
            if cut < per_line*0.45:
                # no clean separator: fall back to the last position outside quotes
                depth_s=depth_d=0; last=-1
                for k,ch in enumerate(window):
                    if ch=="'" and depth_d==0: depth_s^=1
                    elif ch=='"' and depth_s==0: depth_d^=1
                    if depth_s==0 and depth_d==0: last=k
                cut = last if last>per_line*0.35 else per_line-1
            lines.append(w[:cut+1]); w=w[cut+1:]
        cur=w
        if len(lines)>=max_lines: break
    if cur and len(lines)<max_lines: lines.append(cur)
    if not lines: lines=[cmd[:per_line]]
    if len(lines)>max_lines:
        lines=lines[:max_lines]
    # if content was dropped, mark the last line
    if len(" ".join(lines))<len(cmd)-1:
        lines[-1]=lines[-1][:per_line-2].rstrip()+" …"
    return lines

def _fit_title(title,size=29):
    """Shrink long titles so they never wrap into the hairline rule below."""
    n=len(title)
    if n<=52: return size
    if n<=66: return 25
    if n<=82: return 22
    return 20

def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.88),Inches(11.9),Inches(0.78),
        [[(title,_fit_title(title),INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — Gemini Agent ADK badge, else text fallback
    badge=_logo("comptia-dataplus-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.35),Inches(0.6),width=Inches(2.2))
    else:
        rect(s,Inches(10.75),Inches(0.66),Inches(1.85),Inches(1.12),BLUE)
        txt(s,Inches(10.75),Inches(0.78),Inches(1.85),Inches(0.52),[[("CompTIA",18,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(10.75),Inches(1.32),Inches(1.85),Inches(0.4),[[("DATA+  DA0-001",9,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def process_map(title,stages,kicker=None,color=BLUE,synthesis=None,animate=True):
    """STAGED PROCESS MAP — real rounded-rect stages joined by real connectors with
    arrowheads, each stage numbered, optional synthesis band. stages: list of
    (label, detail). This replaces flow_h wherever a genuine process is taught."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(stages); X0=Inches(0.85); TOTW=Inches(11.63)
    gap=Inches(0.42); cw=int((TOTW-gap*(n-1))/n)
    y=Inches(2.35); ch=Inches(2.35) if synthesis else Inches(3.0)
    spids=[]
    for i,st in enumerate(stages):
        lbl,detail=(st if isinstance(st,tuple) else (st,""))
        x=int(X0+(cw+gap)*i)
        box=roundrect(s,x,y,cw,ch,LIGHT,line=LINE)
        rect(s,x,y,cw,Inches(0.11),color)
        bd=Inches(0.62)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,[[(str(i+1),24,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        # label sits in its own fixed band; the caption gets a single line beneath it
        lbl=_ellipsis(lbl,44)
        lsz=13 if len(lbl)<=28 else (12 if len(lbl)<=38 else 11)
        txt(s,x+Inches(0.14),int(y+Inches(1.06)),cw-Inches(0.28),Inches(0.94),
            [[(lbl,lsz,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
        if detail:
            # ONE line only, sized to the card — a 2-line caption overflows the card bottom
            det=_ellipsis(detail,34)
            txt(s,x+Inches(0.1),int(y+ch-Inches(0.52)),cw-Inches(0.2),Inches(0.38),
                [[(det,9,GREY,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
        spids.append(box.shape_id)
        if i<n-1:   # REAL connector in the gap, with an arrowhead
            cy=int(y+ch/2)
            cn=connector(s,int(x+cw+Inches(0.06)),cy,int(x+cw+gap-Inches(0.06)),cy,color)
            spids.append(cn.shape_id)
    if synthesis:
        by=int(y+ch+Inches(0.34))
        rect(s,Inches(0.85),by,Inches(11.63),Inches(1.15),LIGHT)
        rect(s,Inches(0.85),by,Inches(0.11),Inches(1.15),color)
        txt(s,Inches(1.15),int(by+Inches(0.12)),Inches(11.1),Inches(0.32),
            [[(synthesis[0].upper(),11,color,True)]])
        txt(s,Inches(1.15),int(by+Inches(0.46)),Inches(11.1),Inches(0.62),
            [[(synthesis[1],13,INK,False)]])
    if animate: _build_on_click(s,spids)
    footer(s); return s

def decision_map(title,question,yes,no,kicker=None,color=VIOLET,note=None):
    """A real decision diamond with two branches drawn as connectors — used for
    'which pattern do I choose' teaching moments."""
    s=head(slide(),title,kicker,kcolor=color)
    # a diamond's usable text area is ~50% of its box — size generously or text spills
    # past the facets. 4.6 x 2.7 fits 2-3 short lines at 12pt.
    dx,dy,dw,dh=Inches(0.85),Inches(2.75),Inches(4.6),Inches(2.7)
    d=diamond(s,dx,dy,dw,dh,color); label_in(d,question,12,WHITE)
    bx=Inches(6.1); bw=Inches(6.35); bh=Inches(1.5)
    ys=[Inches(2.15),Inches(4.35)]
    for (hdr,items),by,col in zip([yes,no],ys,[TEAL,AMBER]):
        b=roundrect(s,bx,by,bw,bh,LIGHT,line=LINE)
        rect(s,bx,by,Inches(0.11),bh,col)
        txt(s,bx+Inches(0.3),int(by+Inches(0.16)),bw-Inches(0.55),Inches(0.4),[[(hdr,15,col,True)]])
        txt(s,bx+Inches(0.3),int(by+Inches(0.6)),bw-Inches(0.55),Inches(0.82),[[(items,12,INK,False)]])
        connector(s,int(dx+dw),int(dy+dh/2),bx,int(by+bh/2),col)
    if note:
        rect(s,Inches(0.85),Inches(6.15),Inches(11.63),Inches(0.72),LIGHT)
        txt(s,Inches(1.15),Inches(6.28),Inches(11.1),Inches(0.5),[[(note,12,GREY,False)]])
    footer(s); return s

def compare_table(title,headers,rows,kicker=None,accent=BLUE,note=None):
    """A real comparison matrix — the substantive alternative to two bullet columns."""
    s=head(slide(),title,kicker,kcolor=accent)
    X0=Inches(0.85); TOTW=Inches(11.63); ncol=len(headers)
    first=int(TOTW*0.26); rest=int((TOTW-first)/(ncol-1))
    widths=[first]+[rest]*(ncol-1)
    y=Inches(1.95); hh=Inches(0.52)
    x=X0
    for i,h in enumerate(headers):
        col=accent if i==0 else PALETTE[(i-1)%len(PALETTE)]
        rect(s,x,y,widths[i],hh,col)
        txt(s,x+Inches(0.14),y,widths[i]-Inches(0.28),hh,[[(h,13,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
        x+=widths[i]
    # body must end above the note band AND the footer (7.05). Budget it explicitly.
    ry=int(y+hh)
    BOTTOM=Inches(6.88)-(Inches(1.24) if note else Emu(0))
    avail=BOTTOM-ry
    rh=int(min(Inches(0.78),avail/max(len(rows),1)))
    for r,row in enumerate(rows):
        x=X0
        for i,cell in enumerate(row):
            fill=LIGHT if r%2==0 else WHITE
            rect(s,x,ry,widths[i],rh,fill,line=LINE)
            bold=(i==0)
            txt(s,x+Inches(0.14),ry,widths[i]-Inches(0.28),rh,
                [[(cell,11.5,INK if bold else GREY,bold)]],anchor=MSO_ANCHOR.MIDDLE)
            x+=widths[i]
        ry+=rh
    if note:
        rect(s,X0,int(ry+Inches(0.22)),TOTW,Inches(0.92),LIGHT)
        rect(s,X0,int(ry+Inches(0.22)),Inches(0.11),Inches(0.92),accent)
        txt(s,X0+Inches(0.3),int(ry+Inches(0.34)),TOTW-Inches(0.6),Inches(0.66),
            [[("WHEN IT MATTERS  ",11,accent,True),(note,12,INK,False)]])
    footer(s); return s

def worked_example(title,intro,code,explain,kicker=None,accent=TEAL):
    """A worked example: the code on the left, the line-by-line reading on the right.
    This is what turns a decorative lab slide into a teaching slide."""
    s=head(slide(),title,kicker,kcolor=accent)
    txt(s,Inches(0.85),Inches(1.9),Inches(11.63),Inches(0.46),[[(intro,15,GREY,False)]])
    cx,cw=Inches(0.85),Inches(6.5)
    rect(s,cx,Inches(2.5),cw,Inches(4.15),RGBColor(0x0B,0x12,0x20))
    tb=s.shapes.add_textbox(cx+Inches(0.22),Inches(2.62),cw-Inches(0.44),Inches(3.9))
    tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(code):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(2)
        r=p.add_run(); r.text=ln; r.font.size=Pt(11); r.font.name="Consolas"
        col=RGBColor(0x9C,0xDC,0xFE)
        st=ln.strip()
        if st.startswith("#"): col=RGBColor(0x6A,0x99,0x55)
        elif "=" in ln and not st.startswith(("def","class")): col=RGBColor(0xD4,0xD4,0xD4)
        if st.startswith(("def ","class ","from ","import ")): col=RGBColor(0xC5,0x86,0xC0)
        r.font.color.rgb=col
    ex,ew=Inches(7.65),Inches(4.83)
    for i,(lbl,body) in enumerate(explain[:4]):
        y=int(Inches(2.5)+(Inches(1.0)+Inches(0.05))*i)
        col=PALETTE[i%len(PALETTE)]
        rect(s,ex,y,ew,Inches(1.0),LIGHT); rect(s,ex,y,Inches(0.09),Inches(1.0),col)
        txt(s,ex+Inches(0.26),int(y+Inches(0.1)),ew-Inches(0.45),Inches(0.32),[[(lbl,12,col,True)]])
        txt(s,ex+Inches(0.26),int(y+Inches(0.42)),ew-Inches(0.45),Inches(0.52),[[(body,11,INK,False)]])
    footer(s); return s

def steps_slide(act_title,steps,kicker,accent=TEAL,part=None,start=1):
    """Substantive lab procedure: numbered steps WITH their commands, up to 5 per
    slide. Replaces the old one-sentence step slide."""
    s=head(slide(),act_title+(f" — {part}" if part else ""),kicker,kcolor=accent)
    y0=Inches(1.92); n=len(steps); gapy=Inches(0.1)
    AVAIL=Inches(4.92)          # 1.92 → 6.84, clear of the 7.05 footer
    # a step whose command spans several lines needs a taller card
    _lines=[max(1,len([l for l in (c or "").split("\n") if l.strip()])) for _,c in steps]
    _units=sum(min(x,5) for x in _lines) or n
    rh_max=int(min(Inches(1.55),(AVAIL-gapy*(n-1))/max(n,1)*(5.0/max(_units/max(n,1),1.0))))
    rh=int(min(Inches(1.12) if max(_lines)<=1 else Inches(1.55),
               (AVAIL-gapy*(n-1))/max(n,1)))
    for i,(text,cmd) in enumerate(steps):
        y=int(y0+(rh+gapy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,Inches(0.85),y,Inches(11.63),rh,LIGHT); rect(s,Inches(0.85),y,Inches(0.09),rh,col)
        bd=Inches(0.4)
        oval(s,Inches(1.06),int(y+rh/2-bd/2),bd,bd,col)
        txt(s,Inches(1.06),int(y+rh/2-bd/2),bd,bd,[[(str(start+i),13,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        if cmd:
            # one line only in the title band — a 2-line label is overlapped by the code block
            lab=_ellipsis(text,190)
            lsz=12 if len(lab)<=104 else 10.5
            txt(s,Inches(1.62),int(y+Inches(0.04)),Inches(10.6),Inches(0.40),[[(lab,lsz,INK,True)]])
            rect(s,Inches(1.62),int(y+Inches(0.44)),Inches(10.5),int(rh-Inches(0.54)),RGBColor(0x0B,0x12,0x20))
            # WRAP the command over up to 3 lines rather than slicing it into an
            # unrunnable fragment — a partial command teaches the wrong thing.
            raw=str(cmd)
            # A command the author already split across lines is pre-formatted: render it
            # verbatim so no line is re-broken inside a URL or a quoted literal.
            _pre=[l for l in str(cmd).split("\n") if l.strip()]
            _is_pre = len(_pre)>1 and "<<'EOF'" not in raw and '<<"EOF"' not in raw
            def _split_py_c(t):
                """Split `python3 -c "a;b;c"` at its statement semicolons so no line
                ever breaks inside a regex or string literal."""
                m3=re.match(r'^(.*?-c\s+)([\'"])(.*)\2\s*$', t, re.S)
                if not m3: return None
                pre,q,prog=m3.group(1),m3.group(2),m3.group(3)
                parts=[]; cur=""; ds=dd=0
                for ch in prog:
                    if ch=="'" and dd==0: ds^=1
                    elif ch=='"' and ds==0: dd^=1
                    cur+=ch
                    if ch==";" and ds==0 and dd==0:
                        parts.append(cur); cur=""
                if cur: parts.append(cur)
                if len(parts)<2: return None
                out=[pre+q+parts[0]]
                out+= ["    "+x for x in parts[1:-1]]
                out.append("    "+parts[-1]+q)
                return out
            # A multi-line heredoc script belongs in the Learner Guide / lab file, not on a
            # slide. Showing the first fragment would teach an unrunnable command, so show
            # the invocation and point at the file that carries the whole script.
            if "<<'EOF'" in raw or '<<"EOF"' in raw or raw.count("\n") >= 4:
                head_line=raw.split("\n")[0].strip()
                fname=""
                m2=re.search(r"cat\s*>\s*([\w./-]+)", head_line)
                if m2: fname=m2.group(1)
                flat=(f"{head_line}   # full script in the Learner Guide / lab README"
                      if fname else head_line)
            else:
                flat=" ".join(raw.split())
            if _is_pre:
                wrapped=_pre[:6]
            else:
                pysplit=_split_py_c(flat) if len(flat)>110 else None
                if pysplit and len(pysplit)<=6:
                    wrapped=[_ellipsis(x,132) for x in pysplit]
                else:
                    wrapped=_wrap_cmd(flat,per_line=110,max_lines=6)
            _mx=max((len(x) for x in wrapped), default=0)
            csize={1:10.5,2:9.0,3:7.8}.get(len(wrapped),6.4)
            if _mx>100: csize=min(csize,6.4)
            elif _mx>84: csize=min(csize,7.4)
            runs=[[("$ "+wrapped[0],csize,RGBColor(0x9C,0xDC,0xFE),False)]]
            runs+=[[("   "+ln,csize,RGBColor(0x9C,0xDC,0xFE),False)] for ln in wrapped[1:]]
            txt(s,Inches(1.78),int(y+Inches(0.44)),Inches(10.2),int(rh-Inches(0.54)),
                runs,anchor=MSO_ANCHOR.MIDDLE,space=0,font="Consolas",wrap=False)
        else:
            txt(s,Inches(1.62),y,Inches(10.6),rh,[[(text,12.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s

def chart_slide(title,categories,series,kicker=None,accent=BLUE,
                kind="bar",insight=None,number_format='0'):
    """NATIVE PowerPoint chart (fully editable, not a picture) + an insight band.
    kind: bar | column | line | pie | doughnut. series: list of (name, values)."""
    s=head(slide(),title,kicker,kcolor=accent)
    cd=ChartData(); cd.categories=categories
    for nm,vals in series: cd.add_series(nm,vals,number_format)
    ctype={"bar":XL_CHART_TYPE.BAR_CLUSTERED,"column":XL_CHART_TYPE.COLUMN_CLUSTERED,
           "line":XL_CHART_TYPE.LINE_MARKERS,"pie":XL_CHART_TYPE.PIE,
           "doughnut":XL_CHART_TYPE.DOUGHNUT}.get(kind,XL_CHART_TYPE.COLUMN_CLUSTERED)
    ch_h=Inches(3.55) if insight else Inches(4.75)
    gf=s.shapes.add_chart(ctype,Inches(0.85),Inches(1.95),Inches(11.63),ch_h,cd)
    ch=gf.chart
    ch.has_title=False
    ch.font.size=Pt(12); ch.font.name="Arial"; ch.font.color.rgb=INK
    if kind in ("pie","doughnut") or len(series)>1:
        ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout=False; ch.legend.font.size=Pt(11)
    else:
        ch.has_legend=False
    # house palette per point (pie/doughnut) or per series (bar/column/line)
    try:
        if kind in ("pie","doughnut"):
            pts=ch.plots[0]
            for i,pt in enumerate(pts.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb=PALETTE[i%len(PALETTE)]
        else:
            for i,sr in enumerate(ch.series):
                col=PALETTE[i%len(PALETTE)]
                if kind=="line":
                    sr.format.line.color.rgb=col; sr.format.line.width=Pt(2.5)
                else:
                    sr.format.fill.solid(); sr.format.fill.fore_color.rgb=col
    except Exception:
        pass
    try:
        pl=ch.plots[0]; pl.has_data_labels=True
        dl=pl.data_labels; dl.font.size=Pt(10); dl.font.color.rgb=INK
        dl.number_format=number_format; dl.number_format_is_linked=False
        if kind in ("pie","doughnut"): dl.position=XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    if insight:
        by=Inches(5.72)
        rect(s,Inches(0.85),by,Inches(11.63),Inches(1.1),LIGHT)
        rect(s,Inches(0.85),by,Inches(0.11),Inches(1.1),accent)
        txt(s,Inches(1.15),int(by+Inches(0.12)),Inches(11.1),Inches(0.3),
            [[("WHAT THE DATA SHOWS",11,accent,True)]])
        txt(s,Inches(1.15),int(by+Inches(0.44)),Inches(11.1),Inches(0.6),
            [[(insight,12,INK,False)]])
    footer(s); return s

def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker,objective=None,test=None):
    """Lab briefing — now a full teaching slide: the tag chip, the description, and a
    3-tile band covering objective / deliverable / toolchain, plus the success test."""
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.88),Inches(1.7),Inches(0.46),TEAL)
    txt(s,Inches(0.85),Inches(1.88),Inches(1.7),Inches(0.46),[[(tag,15,WHITE,True)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if objective:
        txt(s,Inches(2.72),Inches(1.88),Inches(9.7),Inches(0.46),
            [[(objective,12,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(0.85),Inches(2.5),Inches(11.63),Inches(1.15),[[(desc,17,INK,False)]])
    tiles=[(BLUE,"YOU'LL BUILD",build),(TEAL,"TOOLCHAIN",services),
           (VIOLET,"DONE WHEN",test or "The lab runs end to end without error.")]
    tw=Inches(3.71); xs=[Inches(0.85),Inches(4.81),Inches(8.77)]
    for (col,lbl,body),x in zip(tiles,xs):
        rect(s,x,Inches(3.8),tw,Inches(2.05),LIGHT); rect(s,x,Inches(3.8),tw,Inches(0.1),col)
        txt(s,x+Inches(0.24),Inches(3.98),tw-Inches(0.45),Inches(0.34),[[(lbl,11,col,True)]])
        txt(s,x+Inches(0.24),Inches(4.34),tw-Inches(0.45),Inches(1.4),[[(body,12,INK,False)]])
    footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],
            anchor=MSO_ANCHOR.MIDDLE,font="Consolas",wrap=False)
    footer(s); return s
def test_slide(act_title,text,kicker,troubleshoot=None):
    """Verification slide — the success criterion PLUS a troubleshooting band, so it
    teaches diagnosis rather than stating one sentence."""
    s=head(slide(),act_title,kicker,TEAL)
    GREEN=RGBColor(0x12,0x7A,0x3E)
    rect(s,Inches(0.85),Inches(1.95),Inches(11.63),Inches(2.15),RGBColor(0xE8,0xF7,0xEE))
    rect(s,Inches(0.85),Inches(1.95),Inches(0.11),Inches(2.15),GREEN)
    txt(s,Inches(1.2),Inches(2.12),Inches(11),Inches(0.44),[[("✅  Expected result",17,GREEN,True)]])
    txt(s,Inches(1.2),Inches(2.62),Inches(11.0),Inches(1.35),[[(text,15,INK,False)]])
    tb=troubleshoot or [
        ("Nothing happens","Check the .env file is in the labs folder and the key has no quotes or spaces."),
        ("Auth or 401 error","Re-copy the API key from AI Studio; confirm GOOGLE_GENAI_USE_VERTEXAI=0."),
        ("ModuleNotFoundError","Run uv sync again, and prefix commands with uv run so the venv is used."),
    ]
    txt(s,Inches(0.85),Inches(4.32),Inches(11.63),Inches(0.34),
        [[("IF IT DOESN'T WORK",11,AMBER,True)]])
    tw=Inches(3.71); xs=[Inches(0.85),Inches(4.81),Inches(8.77)]
    for i,(sym,fix) in enumerate(tb[:3]):
        x=xs[i]
        rect(s,x,Inches(4.7),tw,Inches(1.72),LIGHT); rect(s,x,Inches(4.7),tw,Inches(0.09),AMBER)
        txt(s,x+Inches(0.24),Inches(4.87),tw-Inches(0.45),Inches(0.36),[[(sym,12.5,INK,True)]])
        txt(s,x+Inches(0.24),Inches(5.26),tw-Inches(0.45),Inches(1.05),[[(fix,11,GREY,False)]])
    footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

